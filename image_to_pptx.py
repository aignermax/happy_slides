import os
from pptx import Presentation
from pptx.util import Inches
from natsort import natsorted
from PIL import Image, ExifTags
import zipfile
from lxml import etree

# Description:
# This script processes all images located in the script's directory and its subdirectories.
# It creates a PowerPoint presentation where each image is placed on a separate slide.
# Images are correctly oriented based on their EXIF data and are inserted in the order of their folders and filenames.

# Path to the main folder containing subdirectories with images
main_folder = r"."

# Set the working directory
os.chdir(main_folder)
print(f"Main folder set to: {main_folder}")
print(f"Working directory set to: {os.getcwd()}")

# Create a PowerPoint presentation object
presentation = Presentation()

# Function to read EXIF data and correct image orientation
def fix_image_orientation(image_path):
    try:
        image = Image.open(image_path)
        for orientation in ExifTags.TAGS.keys():
            if ExifTags.TAGS[orientation] == 'Orientation':
                break
        exif = image._getexif()
        if exif is not None:
            orientation = exif.get(orientation, 1)
            if orientation == 3:  # Rotate 180 degrees
                image = image.rotate(180, expand=True)
            elif orientation == 6:  # Rotate 270 degrees
                image = image.rotate(270, expand=True)
            elif orientation == 8:  # Rotate 90 degrees
                image = image.rotate(90, expand=True)
        return image
    except Exception as e:
        print(f"Error processing {image_path}: {e}")
        return None

# Function to collect and sort images from a folder (including subfolders)
def get_images_sorted(folder):
    images = []
    for root, _, files in os.walk(folder):
        for file in files:
            if file.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
                images.append(os.path.join(root, file))
    # Sort images in the current folder
    return natsorted(images)

# Sort main folders alphabetically
main_folders = [os.path.join(main_folder, d) for d in os.listdir(main_folder) if os.path.isdir(os.path.join(main_folder, d))]
main_folders = natsorted(main_folders)

# Process images from each main folder in the correct order
for folder in main_folders:
    print(f"Processing folder: {folder}")
    images = get_images_sorted(folder)
    for image_path in images:
        # Load image with correct orientation
        fixed_image = fix_image_orientation(image_path)
        if fixed_image is None:
            continue
        
        # Save temporary image as python-pptx works with file paths
        temp_path = "temp_image.jpg"
        fixed_image.save(temp_path)

        # Add a new slide for the image
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        # Add image and adjust size proportionally
        try:
            picture = slide.shapes.add_picture(temp_path, Inches(0), Inches(0), width=slide_width)
            picture.left = (slide_width - picture.width) // 2
            picture.top = (slide_height - picture.height) // 2
        except Exception as e:
            print(f"Error adding image: {image_path}. Error: {e}")

# Delete temporary file
if os.path.exists("temp_image.jpg"):
    os.remove("temp_image.jpg")

# Save the presentation
output_file = "Photos_Presentation.pptx"
presentation.save(output_file)
print(f"Presentation saved: {os.path.join(main_folder, output_file)}")

# Function to add slide transitions
def add_slide_transitions(pptx_file, duration=4):
    with zipfile.ZipFile(pptx_file, "r") as pptx_zip:
        # Extract the file into a temporary folder
        temp_folder = "temp_pptx"
        pptx_zip.extractall(temp_folder)

    # Find the presentation XML
    slide_dir = os.path.join(temp_folder, "ppt", "slides")
    slide_files = [f for f in os.listdir(slide_dir) if f.endswith(".xml")]

    # Add transitions to each slide
    for slide_file in slide_files:
        slide_path = os.path.join(slide_dir, slide_file)
        tree = etree.parse(slide_path)
        root = tree.getroot()
        ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        
        # Search for the <p:sld> element
        slide_element = root.find(".//p:sld", namespaces=ns)
        if slide_element is None:
            print(f"Warning: No <p:sld> element found in file {slide_file}. Skipping this slide.")
            continue

        # Add transition element
        transition = etree.SubElement(
            slide_element,
            "{http://schemas.openxmlformats.org/presentationml/2006/main}transition",
            attrib={"transition": "fade"}
        )
        advance_on_time = etree.SubElement(
            transition, "{http://schemas.openxmlformats.org/presentationml/2006/main}advanceOnTime"
        )
        advance_on_time.text = "1"
        duration_elem = etree.SubElement(
            transition, "{http://schemas.openxmlformats.org/presentationml/2006/main}duration"
        )
        duration_elem.text = str(duration)

        # Save the updated slide
        with open(slide_path, "wb") as f:
            f.write(etree.tostring(tree, pretty_print=True))

    # Recreate the ZIP file
    new_pptx_file = "Photos_Presentation_Transitions.pptx"
    with zipfile.ZipFile(new_pptx_file, "w") as pptx_zip:
        for root, _, files in os.walk(temp_folder):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, temp_folder)
                pptx_zip.write(file_path, arcname)

    # Delete temporary files
    for root, dirs, files in os.walk(temp_folder, topdown=False):
        for name in files:
            os.remove(os.path.join(root, name))
        for name in dirs:
            os.rmdir(os.path.join(root, name))
    os.rmdir(temp_folder)

    print(f"New presentation with transitions saved: {new_pptx_file}")

# Add transitions
add_slide_transitions(output_file, duration=4)
